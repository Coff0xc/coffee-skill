from __future__ import annotations

import csv
import shutil
import textwrap
import zipfile
from copy import deepcopy
from datetime import datetime
from decimal import Decimal
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_BREAK, WD_PARAGRAPH_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
QUALITY_ROOT = ROOT / "evals" / "quality"
GOLDEN_ROOT = QUALITY_ROOT / "golden-responses"
CASES_ROOT = QUALITY_ROOT / "cases"
WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
DOCX_NS = {
    "w": WORD_NS,
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def reset_dir(path: Path) -> None:
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True, exist_ok=True)


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(textwrap.dedent(text).strip() + "\n", encoding="utf-8")


def create_png(path: Path, size: tuple[int, int], title: str, lines: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", size, "#F7F4ED")
    draw = ImageDraw.Draw(image)
    draw.rectangle([0, 0, size[0], 72], fill="#12343B")
    draw.text((28, 24), title, fill="#FFFFFF")
    y = 104
    colors = ["#D9A441", "#2F7F95", "#D45D4C", "#6B7280", "#234E52"]
    for idx, line in enumerate(lines):
        x = 32 + (idx % 3) * (size[0] // 3)
        if idx and idx % 3 == 0:
            y += 150
        draw.rectangle([x, y, x + size[0] // 3 - 58, y + 92], outline=colors[idx % len(colors)], width=3)
        draw.text((x + 16, y + 20), line[:42], fill="#172026")
    image.save(path)


def docx_xml_text(root: ET.Element) -> str:
    values = [
        elem.text
        for elem in root.iter()
        if elem.tag in {qn("w:t"), qn("w:delText"), qn("w:instrText")} and elem.text
    ]
    return " ".join(values)


def build_ui_golden() -> None:
    case = GOLDEN_ROOT / "ui-admin-dashboard-visual-gate"
    reset_dir(case)
    write_text(
        case / "output" / "index.html",
        """
        <!doctype html>
        <html lang="en">
        <head>
          <meta charset="utf-8" />
          <meta name="viewport" content="width=device-width, initial-scale=1" />
          <title>Billing Operations Console</title>
          <style>
            :root {
              --surface: #f8faf7;
              --ink: #182225;
              --muted: #5f6b6f;
              --accent: #b75f3a;
              --ok: #2f7f5f;
              --warn: #b7791f;
              --line: #d9ded7;
            }
            * { box-sizing: border-box; }
            body { margin: 0; font-family: Arial, sans-serif; color: var(--ink); background: var(--surface); }
            main { max-width: 1160px; margin: 0 auto; padding: 24px; }
            .toolbar { display: flex; align-items: center; justify-content: space-between; gap: 16px; }
            button { border: 1px solid var(--line); background: #fff; color: var(--ink); padding: 8px 12px; border-radius: 6px; }
            button:focus { outline: 3px solid #87b5c4; outline-offset: 2px; }
            button:disabled, .disabled { opacity: .45; cursor: not-allowed; }
            .grid { display: grid; grid-template-columns: 1.2fr .8fr; gap: 18px; margin-top: 18px; }
            .panel { border: 1px solid var(--line); background: #fff; padding: 18px; border-radius: 8px; }
            .state { display: grid; gap: 10px; }
            .loading, .empty, .error { padding: 10px 12px; border-left: 4px solid var(--accent); background: #fbf6f2; }
            .metric { font-size: 32px; font-weight: 700; }
            @media (max-width: 720px) {
              main { padding: 16px; }
              .toolbar, .grid { display: block; }
              .panel { margin-top: 14px; }
            }
          </style>
        </head>
        <body>
          <main>
            <section class="toolbar">
              <div>
                <h1>Billing Operations Console</h1>
                <p>Exception review queue, owner triage, and state coverage.</p>
              </div>
              <button type="button">Refresh queue</button>
            </section>
            <section class="grid">
              <article class="panel">
                <h2>Exception workload</h2>
                <p class="metric">271</p>
                <p>Open exceptions with West concentration and tax-code review pressure.</p>
                <button type="button" class="disabled" disabled>Assign selected</button>
              </article>
              <aside class="panel state">
                <h2>States</h2>
                <div class="loading">Loading: fetching latest queue status.</div>
                <div class="empty">Empty: no exceptions match the current filter.</div>
                <div class="error">Error: source workbook needs review before refresh.</div>
              </aside>
            </section>
          </main>
        </body>
        </html>
        """,
    )
    write_text(
        case / "evaluation-notes.md",
        """
        # UI Need Package

        Product type: SaaS/admin operations console for repeated billing queue triage.
        Design system: semantic tokens, restrained panels, clear action hierarchy, and domain labels.
        State coverage: loading, empty, error, disabled, focus, and responsive behavior are represented.
        Accessibility: focus rings, semantic main/section/article/aside structure, and native button controls.
        Browser: desktop and mobile PNG evidence is included for release-gate scoring.
        """,
    )
    create_png(case / "screenshots" / "desktop.png", (1280, 800), "Desktop dashboard render", ["Queue panel", "State rail", "Responsive table"])
    create_png(case / "screenshots" / "mobile.png", (390, 844), "Mobile dashboard render", ["Toolbar", "Queue", "States"])


def build_dev_golden() -> None:
    case = GOLDEN_ROOT / "dev-repo-repair-ci-gate"
    reset_dir(case)
    write_text(
        case / "src" / "billing.py",
        """
        from __future__ import annotations

        from decimal import Decimal, ROUND_HALF_UP


        def normalize_amount(value: str | int | float | Decimal) -> Decimal:
            if isinstance(value, Decimal):
                return value
            normalized = str(value).strip().replace("$", "").replace(",", "")
            return Decimal(normalized)


        def invoice_total(lines, discount_rate=Decimal("0"), tax_rate=Decimal("0")) -> Decimal:
            subtotal = Decimal("0")
            for line in lines:
                quantity = normalize_amount(line["quantity"])
                unit_price = normalize_amount(line["unit_price"])
                subtotal += quantity * unit_price
            discounted = subtotal * (Decimal("1") - Decimal(str(discount_rate)))
            taxed = discounted * (Decimal("1") + Decimal(str(tax_rate)))
            return taxed.quantize(Decimal("0.01"), rounding=ROUND_HALF_UP)
        """,
    )
    shutil.copyfile(
        CASES_ROOT / "dev-repo-repair-ci-gate" / "input" / "requirements.lock",
        case / "requirements.lock",
    )
    write_text(
        case / "repair-notes.md",
        """
        # Repair Notes

        Need Package: failing pytest path, CI log, local AGENTS rules, and billing module contract.
        Root cause: currency strings with commas were not normalized before Decimal conversion.
        Fast inner loop: ran the targeted invoice math behavior before broader checks.
        CI: the fixture CI log points to pytest failures in test_billing.py.
        Pytest: the repaired normalize_amount and invoice_total paths satisfy currency parsing, discount_rate, and tax_rate behavior.
        Lockfile: requirements.lock is intentionally unchanged to avoid dependency noise.
        """,
    )


def apply_text_style(run, size: int, bold: bool = False, color: str = "182225") -> None:
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*(int(color[i : i + 2], 16) for i in (0, 2, 4)))


def add_textbox(slide, left, top, width, height, text: str, size: int = 18, bold: bool = False, color: str = "182225"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = PptPt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*(int(color[i : i + 2], 16) for i in (0, 2, 4)))
    return box


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    rgb = RGBColor(*(int(color[i : i + 2], 16) for i in (0, 2, 4)))
    shape.fill.fore_color.rgb = rgb
    shape.line.color.rgb = rgb


def add_source(slide, text: str = "Source: Q2 billing operations outline; internal review notes.") -> None:
    add_textbox(slide, PptInches(0.45), PptInches(6.92), PptInches(12), PptInches(0.22), text, size=7, color="59656A")


def build_ppt_golden() -> None:
    case = GOLDEN_ROOT / "office-ppt-aesthetic"
    reset_dir(case)
    write_text(
        case / "deck-outline.md",
        """
        # Deck Outline

        Slide 1 claim: The operating question is no longer volume alone; the review workload now determines throughput.
        Proof object: title + metric rail for exception volume and review hours.
        Source: Q2 billing operations outline.

        Slide 2 claim: Exceptions improved from April to June, but manual review hours moved the wrong way.
        Proof object: clustered column chart with direct labels.
        Source: outline metrics for April and June.

        Slide 3 claim: Tax-code mismatches are the largest controllable exception family.
        Proof object: ranked mix bars for tax, missing PO, address normalization, and other.
        Source: outline exception mix percentages.

        Slide 4 claim: West region concentration makes local routing rules the next constraint.
        Proof object: region concentration map-style matrix.
        Source: outline region note.

        Slide 5 claim: Earlier validation catches problems but shifts work into review queues.
        Proof object: process flow with workload handoff.
        Source: outline operating note.

        Slide 6 claim: Decisions needed are rule ownership, PO policy, and regional escalation.
        Proof object: decision table.
        Source: leadership review sections.
        """,
    )
    write_text(
        case / "design-system.md",
        """
        # Design System

        Typography: Aptos/Arial sans for labels, large editorial claim titles.
        Palette: off-white canvas, deep ink headings, amber improvement, blue workload, rust risk.
        Chart grammar: direct labels, no heavy legend, bounded proof object per slide.
        Connector grammar: thin straight connectors with semantic direction only.
        Container grammar: hairline panels only when they group a real comparison.
        Footer: quiet Source: note on every evidence slide.
        Banned motifs: generic dashboard, card grid, decorative rounded cards, filler icons, unsupported metrics.
        """,
    )
    write_text(
        case / "contact-sheet-plan.md",
        """
        # Contact Sheet Plan

        Macro-layouts: opening operating question, paired bars, ranked mix proof, regional matrix, process flow, decision table.
        Diversity: six slides use six layout families; no three consecutive slides share the same frame.
        Weak slides: slide 4 needs manual review for regional proof if exact regional figures are added later.
        Repeated slides: none; the contact sheet should read as authored rather than a template pack.
        """,
    )
    write_text(
        case / "comeback-scorecard.md",
        """
        # Comeback Scorecard

        Story: 4 - moves from operating question to decisions.
        Specificity: 5 - uses billing exceptions, manual review hours, tax-code mismatches, and West region.
        Rhythm: 5 - six macro-layouts across six slides.
        Whitespace: 4 - restrained proof objects with open canvas.
        Chart clarity: 4 - direct chart labels and one proof object per analytical slide.
        Typography: 4 - stable claim and label scale.
        Restraint: 5 - no generic card grid or decorative motif.
        Precision: 4 - uses only provided outline metrics.
        Coherence: 5 - consistent palette, footer, and claim structure.
        Render: contact-sheet.png is a deterministic PNG fixture; full PowerPoint renderer is not run in this repository test.
        """,
    )
    create_png(case / "render-checks" / "contact-sheet.png", (1280, 720), "PPT contact sheet evidence", ["Question", "Bars", "Mix", "Region", "Flow", "Decisions"])

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_fill(slide.background.fill, "F7F4ED") if False else None
    add_textbox(slide, PptInches(0.55), PptInches(0.55), PptInches(11.8), PptInches(0.75), "Operating question: why did fewer exceptions create more review work?", 28, True)
    metrics = [("418 -> 271", "invoice exceptions"), ("64 -> 91", "manual review hours / week"), ("38%", "tax-code mismatch share")]
    for idx, (value, label) in enumerate(metrics):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.75 + idx * 4.0), PptInches(2.1), PptInches(3.15), PptInches(1.15))
        set_fill(shape, ["D9A441", "2F7F95", "B75F3A"][idx])
        add_textbox(slide, PptInches(0.95 + idx * 4.0), PptInches(2.26), PptInches(2.7), PptInches(0.38), value, 23, True, "FFFFFF")
        add_textbox(slide, PptInches(0.95 + idx * 4.0), PptInches(2.72), PptInches(2.7), PptInches(0.3), label, 11, False, "FFFFFF")
    add_source(slide)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, PptInches(0.55), PptInches(0.45), PptInches(11.8), PptInches(0.55), "Exceptions improved, but manual review hours moved the wrong way.", 24, True)
    chart_data = CategoryChartData()
    chart_data.categories = ["April", "June"]
    chart_data.add_series("Exceptions", (418, 271))
    chart_data.add_series("Review hours", (64, 91))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, PptInches(1.0), PptInches(1.35), PptInches(7.1), PptInches(4.4), chart_data).chart
    chart.has_legend = True
    add_textbox(slide, PptInches(8.55), PptInches(1.55), PptInches(3.3), PptInches(1.0), "Proof object: volume down, labor up.", 18, True)
    add_textbox(slide, PptInches(8.55), PptInches(2.55), PptInches(3.4), PptInches(1.0), "The validation rules are catching issues earlier but shifting effort to manual review.", 14)
    add_source(slide)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, PptInches(0.55), PptInches(0.45), PptInches(11.8), PptInches(0.55), "Tax-code mismatches are the largest controllable exception family.", 24, True)
    mix = [("Tax-code mismatches", 38, "B75F3A"), ("Missing PO references", 24, "D9A441"), ("Address normalization", 18, "2F7F95"), ("Other", 20, "6B7280")]
    for idx, (label, value, color) in enumerate(mix):
        top = PptInches(1.55 + idx * 0.9)
        add_textbox(slide, PptInches(0.9), top, PptInches(2.5), PptInches(0.3), label, 13)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(3.4), top, PptInches(value / 9), PptInches(0.35))
        set_fill(bar, color)
        add_textbox(slide, PptInches(3.55 + value / 9), top, PptInches(0.8), PptInches(0.3), f"{value}%", 12, True)
    add_source(slide)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, PptInches(0.55), PptInches(0.45), PptInches(11.8), PptInches(0.55), "West region concentration makes local routing rules the next constraint.", 24, True)
    regions = [("West", "highest density", "B75F3A"), ("North", "large-value misses", "D9A441"), ("East", "closed backlog", "2F7F95"), ("South", "data quality gaps", "6B7280")]
    for idx, (region, note, color) in enumerate(regions):
        left = PptInches(0.9 + (idx % 2) * 5.4)
        top = PptInches(1.55 + (idx // 2) * 1.7)
        shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, PptInches(2.1), PptInches(1.05))
        set_fill(shape, color)
        add_textbox(slide, left + PptInches(2.35), top + PptInches(0.05), PptInches(2.25), PptInches(0.32), region, 18, True)
        add_textbox(slide, left + PptInches(2.35), top + PptInches(0.45), PptInches(2.25), PptInches(0.32), note, 12)
    add_source(slide)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, PptInches(0.55), PptInches(0.45), PptInches(11.8), PptInches(0.55), "Earlier validation catches problems but shifts work into review queues.", 24, True)
    steps = [("Validate", "rules catch earlier"), ("Triage", "manual queue grows"), ("Resolve", "ownership gaps surface")]
    for idx, (step, note) in enumerate(steps):
        left = PptInches(0.95 + idx * 4.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, PptInches(2.25), PptInches(2.75), PptInches(1.0))
        set_fill(shape, ["2F7F95", "D9A441", "B75F3A"][idx])
        add_textbox(slide, left + PptInches(0.28), PptInches(2.42), PptInches(2.1), PptInches(0.25), step, 18, True, "FFFFFF")
        add_textbox(slide, left + PptInches(0.28), PptInches(2.78), PptInches(2.1), PptInches(0.25), note, 11, False, "FFFFFF")
        if idx < 2:
            slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + PptInches(2.9), PptInches(2.75), left + PptInches(3.55), PptInches(2.75))
    add_source(slide)

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, PptInches(0.55), PptInches(0.45), PptInches(11.8), PptInches(0.55), "Decisions needed: rule ownership, PO policy, and regional escalation.", 24, True)
    rows = [
        ("Rule ownership", "Finance + Support", "Tax-code mismatches"),
        ("PO policy", "Procurement", "Missing references"),
        ("West escalation", "Regional ops", "Highest density"),
    ]
    table = slide.shapes.add_table(4, 3, PptInches(0.85), PptInches(1.6), PptInches(11.1), PptInches(2.6)).table
    for col, header in enumerate(["Decision", "Owner", "Why now"]):
        table.cell(0, col).text = header
    for row_idx, row in enumerate(rows, 1):
        for col, value in enumerate(row):
            table.cell(row_idx, col).text = value
    add_source(slide)

    (case / "final").mkdir(parents=True, exist_ok=True)
    prs.save(case / "final" / "operating-review.pptx")


def parse_amount(value: str) -> Decimal | None:
    cleaned = value.strip().replace("$", "").replace(",", "")
    if cleaned == "":
        return None
    return Decimal(cleaned)


def normalize_date(value: str) -> str:
    for fmt in ("%Y-%m-%d", "%m/%d/%Y", "%Y/%m/%d"):
        try:
            return datetime.strptime(value, fmt).strftime("%Y-%m-%d")
        except ValueError:
            continue
    return ""


def build_excel_golden() -> None:
    case = GOLDEN_ROOT / "office-excel-parse"
    reset_dir(case)
    write_text(
        case / "parse-audit.md",
        """
        # Parse Audit

        Encoding: UTF-8 fixture read.
        Delimiter: comma CSV with quoted currency values.
        Headers: invoice_id, region, invoice_date, amount, exception_type, review_hours, status.
        Units: amount is USD and review_hours is hours.
        Date parsing: ISO, US slash, and yyyy/mm/dd formats normalized; "not available" becomes a null date.
        Number parsing: dollar signs and comma thousands separators normalized.
        Null handling: blank amount and blank review_hours are preserved and flagged.
        Duplicate rows: invoice INV-1003 appears twice.
        Abnormal values: negative amount -45.00 is flagged.
        Sheet/range inspection: Raw A1:G8 feeds Model A1:K8 and Dashboard A1:D10.
        """,
    )
    write_text(
        case / "workbook-plan.md",
        """
        # Workbook Plan

        Raw: source CSV preserved without overwriting.
        Source: field definitions, source note, and parsing decisions.
        Assumptions: review-hour floor and abnormal-value thresholds.
        Model: normalized date, parsed amount, duplicate flag, null flag, abnormal flag.
        Checks: row count, duplicate count, formula totals, and error scan.
        Dashboard: KPI summary and chart-ready helper range.
        """,
    )
    write_text(
        case / "formula-checks.md",
        """
        # Formula Checks

        Formula-derived values: Dashboard and Checks use formulas rather than pasted prose.
        Bounded range: formulas reference A2:K8 ranges, not full columns.
        Trace: key outputs trace from Raw through Model to Dashboard.
        Error scan: no #REF!, #DIV/0!, #VALUE!, #NAME?, or #N/A values are expected.
        Helper range: Dashboard chart uses a compact helper range for exception mix.
        Chart: bar chart is sourced from Dashboard A12:B15.
        """,
    )
    create_png(case / "render-checks" / "dashboard.png", (1280, 720), "Excel dashboard render", ["KPI block", "Exception mix", "Checks", "Chart"])

    csv_path = CASES_ROOT / "office-excel-parse" / "input" / "messy-billing-exceptions.csv"
    rows = list(csv.DictReader(csv_path.read_text(encoding="utf-8").splitlines()))

    wb = Workbook()
    dashboard = wb.active
    dashboard.title = "Dashboard"
    raw = wb.create_sheet("Raw")
    source = wb.create_sheet("Source")
    assumptions = wb.create_sheet("Assumptions")
    model = wb.create_sheet("Model")
    checks = wb.create_sheet("Checks")

    headers = ["invoice_id", "region", "invoice_date", "amount", "exception_type", "review_hours", "status"]
    raw.append(headers)
    for row in rows:
        raw.append([row[item] for item in headers])
    table = Table(displayName="RawBillingExceptions", ref=f"A1:G{len(rows) + 1}")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
    raw.add_table(table)

    source.append(["Field", "Meaning"])
    for header in headers:
        source.append([header, "Source CSV column"])
    source.append(["Source note", "messy-billing-exceptions.csv"])

    assumptions["A1"] = "Assumption"
    assumptions["B1"] = "Value"
    assumptions["A2"] = "Review hour floor"
    assumptions["B2"] = 0
    assumptions["A3"] = "Negative amount flag"
    assumptions["B3"] = 0

    model_headers = [*headers, "normalized_date", "parsed_amount", "duplicate_flag", "quality_flag"]
    model.append(model_headers)
    seen: set[str] = set()
    total = Decimal("0")
    open_count = 0
    tax_count = 0
    west_count = 0
    duplicate_count = 0
    for row in rows:
        invoice_id = row["invoice_id"]
        amount = parse_amount(row["amount"])
        if amount is not None:
            total += amount
        if row["status"] == "open":
            open_count += 1
        if row["exception_type"] == "Tax Code":
            tax_count += 1
        if row["region"] == "West":
            west_count += 1
        duplicate_flag = "duplicate" if invoice_id in seen else ""
        if duplicate_flag:
            duplicate_count += 1
        seen.add(invoice_id)
        quality = "; ".join(
            item
            for item in [
                "missing-date" if not normalize_date(row["invoice_date"]) else "",
                "missing-amount" if amount is None else "",
                "negative-amount" if amount is not None and amount < 0 else "",
                "missing-hours" if row["review_hours"] == "" else "",
                duplicate_flag,
            ]
            if item
        )
        model.append(
            [
                row["invoice_id"],
                row["region"],
                row["invoice_date"],
                row["amount"],
                row["exception_type"],
                row["review_hours"],
                row["status"],
                normalize_date(row["invoice_date"]),
                float(amount) if amount is not None else None,
                duplicate_flag,
                quality,
            ]
        )

    model_table = Table(displayName="ModelBillingExceptions", ref=f"A1:K{len(rows) + 1}")
    model_table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium4", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
    model.add_table(model_table)

    dashboard["A1"] = "Billing Exceptions Dashboard"
    dashboard["A3"] = "Metric"
    dashboard["B3"] = "Value"
    dashboard["A4"] = "Total parsed amount"
    dashboard["B4"] = "=SUM(Model!I2:I8)"
    dashboard["A5"] = "Open or escalated rows"
    dashboard["B5"] = '=COUNTIFS(Model!G2:G8,"open")+COUNTIFS(Model!G2:G8,"escalated")'
    dashboard["A6"] = "Tax-code rows"
    dashboard["B6"] = '=COUNTIFS(Model!E2:E8,"Tax Code")'
    dashboard["A7"] = "West rows"
    dashboard["B7"] = '=COUNTIFS(Model!B2:B8,"West")'
    dashboard["A8"] = "West tax-code rows"
    dashboard["B8"] = '=COUNTIFS(Model!B2:B8,"West",Model!E2:E8,"Tax Code")'
    dashboard["A9"] = "West tax-code amount"
    dashboard["B9"] = '=SUMIFS(Model!I2:I8,Model!B2:B8,"West",Model!E2:E8,"Tax Code")'
    dashboard["A10"] = "Open amount"
    dashboard["B10"] = '=SUMIFS(Model!I2:I8,Model!G2:G8,"open")'
    helper = [("Tax Code", 3), ("Missing PO", 2), ("Address Normalization", 1), ("Other", 1)]
    dashboard["A11"] = "Exception type"
    dashboard["B11"] = "Rows"
    for idx, (label, value) in enumerate(helper, 12):
        dashboard[f"A{idx}"] = label
        dashboard[f"B{idx}"] = value
    chart = BarChart()
    chart.title = "Exception mix"
    chart.add_data(Reference(dashboard, min_col=2, min_row=11, max_row=15), titles_from_data=True)
    chart.set_categories(Reference(dashboard, min_col=1, min_row=12, max_row=15))
    dashboard.add_chart(chart, "D3")

    checks["A1"] = "Check"
    checks["B1"] = "Value"
    checks["A2"] = "Raw row count"
    checks["B2"] = '=COUNTA(Raw!A2:A8)'
    checks["A3"] = "Duplicate invoices"
    checks["B3"] = '=COUNTIFS(Model!J2:J8,"duplicate")'
    checks["A4"] = "Total amount trace"
    checks["B4"] = "=SUM(Model!I2:I8)"
    checks["A5"] = "Open plus escalated"
    checks["B5"] = '=COUNTIFS(Model!G2:G8,"open")+COUNTIFS(Model!G2:G8,"escalated")'
    checks["A6"] = "West tax amount"
    checks["B6"] = '=SUMIFS(Model!I2:I8,Model!B2:B8,"West",Model!E2:E8,"Tax Code")'

    for sheet in wb.worksheets:
        sheet.freeze_panes = "A2"
        for cell in sheet[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="12343B")
            cell.alignment = Alignment(horizontal="center")
        for column in range(1, sheet.max_column + 1):
            sheet.column_dimensions[chr(64 + column)].width = 18

    (case / "final").mkdir(parents=True, exist_ok=True)
    wb.save(case / "final" / "billing-exceptions.xlsx")


def set_docx_cell_width(cell, width_dxa: int) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_w = tc_pr.first_child_found_in("w:tcW")
    if tc_w is None:
        tc_w = OxmlElement("w:tcW")
        tc_pr.append(tc_w)
    tc_w.set(qn("w:w"), str(width_dxa))
    tc_w.set(qn("w:type"), "dxa")


def add_docx_table_grid(table, widths: list[int]) -> None:
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    tbl_w = OxmlElement("w:tblW")
    tbl_w.set(qn("w:w"), str(sum(widths)))
    tbl_w.set(qn("w:type"), "dxa")
    tbl_pr.append(tbl_w)
    grid = OxmlElement("w:tblGrid")
    for width in widths:
        col = OxmlElement("w:gridCol")
        col.set(qn("w:w"), str(width))
        grid.append(col)
    tbl.insert(1, grid)
    for row in table.rows:
        for idx, cell in enumerate(row.cells):
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            set_docx_cell_width(cell, widths[idx])


def ensure_comment_parts(docx_path: Path) -> None:
    ET.register_namespace("w", WORD_NS)
    ET.register_namespace("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")
    ET.register_namespace("rel", "http://schemas.openxmlformats.org/package/2006/relationships")
    ET.register_namespace("ct", "http://schemas.openxmlformats.org/package/2006/content-types")
    with zipfile.ZipFile(docx_path, "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    document = ET.fromstring(files["word/document.xml"])
    comments = ET.Element(qn("w:comments"))

    body = document.find("w:body", DOCX_NS)
    paragraphs = list(document.findall(".//w:p", DOCX_NS))
    comment_specs = [
        ("0", "Legal needs to confirm indemnity wording.", "Contractual Gaps"),
        ("1", "Evidence link is stale.", "Stale evidence row"),
        ("2", "Reviewer note: evidence is missing for remediation owner.", "Evidence missing"),
    ]
    for comment_id, comment_text, anchor_text in comment_specs:
        comment = ET.SubElement(comments, qn("w:comment"))
        comment.set(qn("w:id"), comment_id)
        comment.set(qn("w:author"), "Quality Eval")
        comment.set(qn("w:date"), "2026-05-16T00:00:00Z")
        p = ET.SubElement(comment, qn("w:p"))
        r = ET.SubElement(p, qn("w:r"))
        t = ET.SubElement(r, qn("w:t"))
        t.text = comment_text

        target_p = next((p for p in paragraphs if anchor_text.lower() in docx_xml_text(p).lower()), paragraphs[0])
        start = ET.Element(qn("w:commentRangeStart"))
        start.set(qn("w:id"), comment_id)
        end = ET.Element(qn("w:commentRangeEnd"))
        end.set(qn("w:id"), comment_id)
        ref_run = ET.Element(qn("w:r"))
        ref = ET.Element(qn("w:commentReference"))
        ref.set(qn("w:id"), comment_id)
        ref_run.append(ref)
        target_p.insert(0, start)
        target_p.append(end)
        target_p.append(ref_run)

    insertion_p = next((p for p in paragraphs if "Executive Summary" in docx_xml_text(p)), paragraphs[0])
    ins = ET.Element(qn("w:ins"))
    ins.set(qn("w:id"), "21")
    ins.set(qn("w:author"), "Quality Eval")
    ins.set(qn("w:date"), "2026-05-16T00:00:00Z")
    ins_run = ET.Element(qn("w:r"))
    ins_text = ET.Element(qn("w:t"))
    ins_text.text = " payment processor stores tokens outside primary region."
    ins_run.append(ins_text)
    ins.append(ins_run)
    insertion_p.append(ins)

    deletion_p = next((p for p in paragraphs if "Required Remediation" in docx_xml_text(p)), paragraphs[-1])
    deletion = ET.Element(qn("w:del"))
    deletion.set(qn("w:id"), "22")
    deletion.set(qn("w:author"), "Quality Eval")
    deletion.set(qn("w:date"), "2026-05-16T00:00:00Z")
    del_run = ET.Element(qn("w:r"))
    del_text = ET.Element(qn("w:delText"))
    del_text.text = " annual review only."
    del_run.append(del_text)
    deletion.append(del_run)
    deletion_p.append(deletion)

    sect_pr = body.find("w:sectPr", DOCX_NS) if body is not None else None
    if sect_pr is not None:
        header_ref = ET.Element(qn("w:headerReference"))
        header_ref.set(qn("w:type"), "default")
        header_ref.set(qn("r:id"), "rId100")
        footer_ref = ET.Element(qn("w:footerReference"))
        footer_ref.set(qn("w:type"), "default")
        footer_ref.set(qn("r:id"), "rId101")
        sect_pr.insert(0, footer_ref)
        sect_pr.insert(0, header_ref)

    header = ET.Element(qn("w:hdr"))
    hp = ET.SubElement(header, qn("w:p"))
    hr = ET.SubElement(hp, qn("w:r"))
    ht = ET.SubElement(hr, qn("w:t"))
    ht.text = "Confidential vendor risk review, Q2 2026"

    footer = ET.Element(qn("w:ftr"))
    fp = ET.SubElement(footer, qn("w:p"))
    fr = ET.SubElement(fp, qn("w:r"))
    ft = ET.SubElement(fr, qn("w:t"))
    ft.text = "Draft owner | Page "
    fld_run = ET.SubElement(fp, qn("w:r"))
    fld_char = ET.SubElement(fld_run, qn("w:fldChar"))
    fld_char.set(qn("w:fldCharType"), "begin")
    instr_run = ET.SubElement(fp, qn("w:r"))
    instr = ET.SubElement(instr_run, qn("w:instrText"))
    instr.text = " PAGE "
    end_run = ET.SubElement(fp, qn("w:r"))
    end = ET.SubElement(end_run, qn("w:fldChar"))
    end.set(qn("w:fldCharType"), "end")

    rels = ET.fromstring(files["word/_rels/document.xml.rels"])
    rel_tag = rels.tag.replace("Relationships", "Relationship")
    for rid, rel_type, target in [
        ("rId99", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments", "comments.xml"),
        ("rId100", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/header", "header1.xml"),
        ("rId101", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/footer", "footer1.xml"),
    ]:
        if not any(rel.attrib.get("Id") == rid for rel in rels):
            rel = ET.SubElement(rels, rel_tag)
            rel.set("Id", rid)
            rel.set("Type", rel_type)
            rel.set("Target", target)

    content = ET.fromstring(files["[Content_Types].xml"])
    override_tag = content.tag.replace("Types", "Override")
    existing_parts = {item.attrib.get("PartName") for item in content}
    for part_name, content_type in [
        ("/word/comments.xml", "application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml"),
        ("/word/header1.xml", "application/vnd.openxmlformats-officedocument.wordprocessingml.header+xml"),
        ("/word/footer1.xml", "application/vnd.openxmlformats-officedocument.wordprocessingml.footer+xml"),
    ]:
        if part_name not in existing_parts:
            override = ET.SubElement(content, override_tag)
            override.set("PartName", part_name)
            override.set("ContentType", content_type)

    files["word/document.xml"] = ET.tostring(document, encoding="utf-8", xml_declaration=True)
    files["word/comments.xml"] = ET.tostring(comments, encoding="utf-8", xml_declaration=True)
    files["word/header1.xml"] = ET.tostring(header, encoding="utf-8", xml_declaration=True)
    files["word/footer1.xml"] = ET.tostring(footer, encoding="utf-8", xml_declaration=True)
    files["word/_rels/document.xml.rels"] = ET.tostring(rels, encoding="utf-8", xml_declaration=True)
    files["[Content_Types].xml"] = ET.tostring(content, encoding="utf-8", xml_declaration=True)

    with zipfile.ZipFile(docx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)


def build_docx_golden() -> None:
    case = GOLDEN_ROOT / "office-docx-format"
    reset_dir(case)
    write_text(
        case / "reading-map.md",
        """
        # Reading Map

        Heading: Title plus Heading 1 and Heading 2 hierarchy.
        Section: Executive Summary, Current Risk Posture, Contractual Gaps, Required Remediation.
        Table: vendor control gaps table with Control, Evidence, Owner, Due Date, Status.
        Comment: existing comments A and B plus one new reviewer comment.
        Tracked: insertion in Executive Summary and deletion in Required Remediation are preserved.
        Header: confidential vendor risk review.
        Footer: draft owner and page number field.
        Field: page number field and TOC/update risk noted.
        Metadata: local fixture author only, no production metadata claim.
        """,
    )
    write_text(
        case / "style-token-map.md",
        """
        # Style Token Map

        Styles: Title, Heading 1, Heading 2, Normal, and table text roles are real Word styles.
        Numbering: real numbering definitions are used for remediation steps.
        Table geometry: tblGrid and tcW widths are explicit.
        Margins: one-inch business memo margins.
        Type scale: 20 pt title, 15 pt Heading 1, 12 pt body.
        Paragraph rhythm: compact spacing after headings and body paragraphs.
        Headers: confidential review label.
        Footers: draft owner plus page field.
        Table gate: only repeated control-gap records use a table.
        """,
    )
    write_text(
        case / "edit-plan.md",
        """
        # Edit Plan

        Preserve original: output is a separate DOCX fixture.
        Minimal: tighten executive summary without changing facts.
        Comment anchor: comments remain anchored to Contractual Gaps and table evidence rows.
        Redline: tracked-change insertion and deletion remain structurally present.
        Tracked-change: do not accept all changes.
        Render: page-1.png and page-2.png are deterministic PNG fixtures; full Word renderer is not run in this repository test.
        """,
    )
    create_png(case / "render-checks" / "page-1.png", (850, 1100), "DOCX page 1 render", ["Header", "Executive Summary", "Risk Posture", "Comments"])
    create_png(case / "render-checks" / "page-2.png", (850, 1100), "DOCX page 2 render", ["Control table", "Remediation", "Footer"])

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    doc.styles["Title"].font.size = Pt(20)
    doc.styles["Heading 1"].font.size = Pt(15)
    doc.styles["Heading 2"].font.size = Pt(13)
    doc.add_heading("Vendor Risk Review", 0)
    doc.add_paragraph("Confidential vendor risk review, Q2 2026.")
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph("The vendor remains usable with remediation, but evidence quality and contractual gaps require owner-level follow-up.")
    doc.add_heading("Current Risk Posture", level=2)
    doc.add_paragraph("Payment processor dependency and evidence freshness are the highest-risk review threads.")
    doc.add_heading("Contractual Gaps", level=2)
    doc.add_paragraph("Indemnity wording needs legal confirmation before renewal. Legal needs to confirm indemnity wording.")
    doc.add_heading("Required Remediation", level=2)
    for item in ["Refresh evidence links", "Confirm token-storage region", "Assign control owners"]:
        doc.add_paragraph(item, style="List Number")

    table = doc.add_table(rows=1, cols=5)
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table.style = "Table Grid"
    headers = ["Control", "Evidence", "Owner", "Due Date", "Status"]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for row in [
        ["Data residency", "Processor attestation", "Security", "2026-06-20", "Open"],
        ["Indemnity", "Contract clause", "Legal", "2026-06-25", "Review"],
        ["Logging", "Stale evidence row", "Vendor", "2026-06-30", "Missing"],
    ]:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value
    add_docx_table_grid(table, [1700, 2300, 1300, 1300, 1200])

    doc.add_section(WD_SECTION.NEW_PAGE)
    doc.add_heading("Appendix: Review Notes", level=1)
    doc.add_paragraph("The table of contents should update after edits.")

    final = case / "final" / "vendor-risk-review.docx"
    final.parent.mkdir(parents=True, exist_ok=True)
    doc.save(final)
    ensure_comment_parts(final)


def main() -> None:
    GOLDEN_ROOT.mkdir(parents=True, exist_ok=True)
    build_ui_golden()
    build_dev_golden()
    build_ppt_golden()
    build_excel_golden()
    build_docx_golden()
    print(f"built golden responses under {GOLDEN_ROOT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
